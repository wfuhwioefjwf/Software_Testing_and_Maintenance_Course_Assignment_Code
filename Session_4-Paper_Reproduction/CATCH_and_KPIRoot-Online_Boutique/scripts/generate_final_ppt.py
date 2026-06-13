from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
TEMPLATE = ROOT / "PPT" / "南开大学模板-50.pptx"
OUTPUT = ROOT / "PPT" / "基于Online Boutique的异常检测与根因定位闭环复现-模板复用版.pptx"

FIG = ROOT / "figures"
IMG_HOME = FIG / "online_boutique_homepage.png"
IMG_CHECKOUT = FIG / "online_boutique_checkout_success.png"
IMG_RUN004 = FIG / "run_004_case_summary.png"
IMG_SUMMARY = FIG / "run_001_to_run_004_summary.png"

SELECTED_TEMPLATE_SLIDES = [1, 2, 3, 4, 5, 11, 13, 18, 23, 24, 25, 31, 35]

PURPLE = RGBColor(0x6E, 0x15, 0x56)
MID_PURPLE = RGBColor(0x8B, 0x1B, 0x6D)
PINK = RGBColor(0xD8, 0x2A, 0xA8)
LIGHT_PURPLE = RGBColor(0xF3, 0xEA, 0xF1)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT = RGBColor(0x24, 0x24, 0x24)
MUTED = RGBColor(0x66, 0x66, 0x66)
LINE = RGBColor(0xDF, 0xD6, 0xE1)

FONT = "Microsoft YaHei"
FONT_BOLD = "Microsoft YaHei UI"


def emu(inches: float):
    return Inches(inches)


def keep_only_selected_slides(prs: Presentation, selected_1_based: list[int]) -> None:
    selected = {idx - 1 for idx in selected_1_based}
    slide_id_list = prs.slides._sldIdLst
    for idx in range(len(slide_id_list) - 1, -1, -1):
        if idx not in selected:
            rel_id = slide_id_list[idx].rId
            prs.part.drop_rel(rel_id)
            del slide_id_list[idx]


def iter_shapes(shapes):
    for shape in shapes:
        yield shape
        if shape.shape_type == 6:
            yield from iter_shapes(shape.shapes)


def clear_text(slide) -> None:
    for shape in iter_shapes(slide.shapes):
        if hasattr(shape, "text_frame"):
            shape.text_frame.clear()


def remove_large_sample_pictures(slide, keep_cover_background=False) -> None:
    for shape in list(slide.shapes):
        if shape.shape_type != 13:
            continue
        if keep_cover_background:
            continue
        element = shape._element
        element.getparent().remove(element)


def set_text_frame(frame, text, size=18, color=TEXT, bold=False, align=PP_ALIGN.LEFT,
                   valign=MSO_ANCHOR.TOP, font=FONT, line_spacing=1.08):
    frame.clear()
    frame.margin_left = emu(0.04)
    frame.margin_right = emu(0.04)
    frame.margin_top = emu(0.02)
    frame.margin_bottom = emu(0.02)
    frame.vertical_anchor = valign
    frame.word_wrap = True
    for idx, line in enumerate(text.split("\n")):
        paragraph = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        paragraph.text = line
        paragraph.alignment = align
        paragraph.line_spacing = line_spacing
        paragraph.space_after = Pt(2)
        for run in paragraph.runs:
            run.font.name = font
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color


def add_text(slide, text, left, top, width, height, size=18, color=TEXT, bold=False,
             align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, font=FONT, line_spacing=1.08):
    box = slide.shapes.add_textbox(left, top, width, height)
    set_text_frame(box.text_frame, text, size, color, bold, align, valign, font, line_spacing)
    return box


def add_rect(slide, left, top, width, height, fill=WHITE, line=None, radius=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
                                   left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    return shape


def add_picture_contain(slide, path, left, top, width, height):
    with Image.open(path) as img:
        img_w, img_h = img.size
    box_ratio = width / height
    img_ratio = img_w / img_h
    if img_ratio >= box_ratio:
        final_w = width
        final_h = width / img_ratio
    else:
        final_h = height
        final_w = height * img_ratio
    pic_left = left + (width - final_w) / 2
    pic_top = top + (height - final_h) / 2
    return slide.shapes.add_picture(str(path), pic_left, pic_top, width=final_w, height=final_h)


def add_title(slide, title, subtitle=None):
    add_text(slide, title, emu(2.52), emu(0.24), emu(3.4), emu(0.34),
             size=13.5, color=PURPLE, bold=True, font=FONT_BOLD)
    if subtitle:
        add_text(slide, subtitle, emu(5.72), emu(0.27), emu(5.2), emu(0.24),
                 size=8.5, color=MUTED)


def add_footer_page_no(slide, num):
    add_text(slide, f"{num:02d}", emu(11.75), emu(6.92), emu(0.45), emu(0.22),
             size=8.5, color=PURPLE, bold=True, align=PP_ALIGN.RIGHT)


def add_bullets(slide, items, left, top, width, height, size=13, color=TEXT, bullet="•"):
    text = "\n".join(f"{bullet} {item}" for item in items)
    return add_text(slide, text, left, top, width, height, size=size, color=color, line_spacing=1.15)


def fill_cover(slide):
    clear_text(slide)
    add_text(slide, "基于 Online Boutique 的\n异常检测与根因定位闭环复现",
             emu(1.12), emu(2.78), emu(10.8), emu(1.12), size=30, color=WHITE,
             bold=True, align=PP_ALIGN.CENTER, font=FONT_BOLD, line_spacing=1.0)
    add_text(slide, "软件工程课程大作业展示",
             emu(4.25), emu(4.04), emu(4.85), emu(0.32), size=14, color=WHITE,
             align=PP_ALIGN.CENTER)
    add_text(slide, "汇报人：王熙康    学号：2312124    日期：2026 年 6 月 6 日",
             emu(2.55), emu(4.92), emu(8.2), emu(0.34), size=12, color=WHITE,
             align=PP_ALIGN.CENTER)
    add_text(slide, "Prometheus / ChaosMesh / Selenium / JMeter / CATCH / KPIRoot",
             emu(3.1), emu(6.43), emu(7.1), emu(0.25), size=8.8, color=RGBColor(0xDD, 0xDD, 0xDD),
             align=PP_ALIGN.CENTER)


def fill_toc(slide):
    clear_text(slide)
    add_text(slide, "目录：", emu(0.88), emu(0.58), emu(1.7), emu(0.52), size=24, color=PINK, bold=True)
    blocks = [
        ("-01-", "项目定位", "课程要求对齐\nOnline Boutique 实验对象"),
        ("-02-", "实验闭环", "监控采集\n故障注入与业务测试"),
        ("-03-", "算法结果", "CATCH 异常检测\nKPIRoot 根因定位"),
        ("-04-", "总结讨论", "工程问题\n个人完成内容与后续工作"),
    ]
    xs = [1.25, 4.0, 6.75, 9.5]
    for x, (num, title, desc) in zip(xs, blocks):
        add_text(slide, num, emu(x), emu(3.02), emu(1.55), emu(0.34), size=14, color=WHITE,
                 bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, title, emu(x - 0.25), emu(3.55), emu(2.08), emu(0.35), size=15, color=WHITE,
                 bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, desc, emu(x - 0.08), emu(4.16), emu(1.75), emu(0.58), size=8.8, color=WHITE,
                 align=PP_ALIGN.CENTER, line_spacing=1.05)


def fill_section(slide):
    clear_text(slide)
    add_text(slide, "01. 项目定位", emu(0.82), emu(2.58), emu(6.8), emu(0.75),
             size=34, color=WHITE, bold=True, font=FONT_BOLD)
    add_text(slide, "从课程四阶段要求出发，构建可展示、可解释的微服务 AIOps 实验闭环。",
             emu(1.0), emu(4.48), emu(5.7), emu(0.58), size=11, color=WHITE, line_spacing=1.08)


def fill_system_intro(slide):
    clear_text(slide)
    add_title(slide, "实验对象：Online Boutique", "复杂开源微服务系统，适合故障注入与根因定位")
    add_picture_contain(slide, IMG_HOME, emu(0.9), emu(1.28), emu(4.25), emu(3.75))
    add_text(slide, "为什么选择它", emu(5.95), emu(1.25), emu(2.5), emu(0.36),
             size=18, color=PURPLE, bold=True)
    add_bullets(slide, [
        "服务数量多，包含 frontend、paymentservice、checkoutservice、recommendationservice 等",
        "业务链路清晰，可构造浏览、加购、结算和推荐访问路径",
        "比简单示例系统更适合展示故障传播与定位过程",
        "本地通过 kubectl port-forward deployment/frontend 8080:8080 访问",
    ], emu(5.95), emu(1.82), emu(5.25), emu(2.45), size=11.2)
    add_text(slide, "课程对应：自行选择更复杂开源微服务系统，并基于自采监控数据完成算法分析。",
             emu(5.95), emu(4.72), emu(5.15), emu(0.55), size=12.2, color=PURPLE, bold=True)
    add_footer_page_no(slide, 4)


def fill_closed_loop(slide):
    clear_text(slide)
    add_title(slide, "总体技术路线", "部署、测试、采集、注入、对齐、检测、定位形成闭环")
    labels = [
        ("Online Boutique\nMinikube 部署", 5.87, 1.25),
        ("Prometheus\nKPI 采集", 9.15, 3.12),
        ("统一宽表\n30s 对齐", 5.87, 5.35),
        ("ChaosMesh\n故障注入", 2.08, 3.12),
    ]
    for text, x, y in labels:
        add_text(slide, text, emu(x), emu(y), emu(2.05), emu(0.48), size=11.5, color=PURPLE,
                 bold=True, align=PP_ALIGN.CENTER, line_spacing=1.0)
    add_text(slide, "CATCH\n何时异常", emu(1.25), emu(5.40), emu(2.2), emu(0.5),
             size=12, color=PURPLE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "KPIRoot\n哪里异常", emu(9.75), emu(5.40), emu(2.2), emu(0.5),
             size=12, color=PURPLE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "关键不是单跑模型，而是把故障标签、业务指标和 Prometheus KPI 统一到同一时间线上。",
             emu(1.15), emu(6.32), emu(10.9), emu(0.32), size=12.5, color=MID_PURPLE,
             bold=True, align=PP_ALIGN.CENTER)
    add_footer_page_no(slide, 5)


def fill_test_evidence(slide):
    clear_text(slide)
    add_title(slide, "Selenium 与 JMeter 测试证据", "功能路径可执行，正式 case 具备稳定业务流量")
    add_picture_contain(slide, IMG_HOME, emu(1.7), emu(1.45), emu(2.38), emu(1.72))
    add_picture_contain(slide, IMG_CHECKOUT, emu(4.78), emu(1.45), emu(2.38), emu(1.72))
    add_picture_contain(slide, IMG_RUN004, emu(7.82), emu(1.45), emu(2.38), emu(1.72))
    add_text(slide, "首页访问", emu(1.54), emu(3.58), emu(2.65), emu(0.58), size=10.5, color=TEXT,
             align=PP_ALIGN.CENTER)
    add_text(slide, "下单成功", emu(4.63), emu(3.58), emu(2.65), emu(0.58), size=10.5, color=TEXT,
             align=PP_ALIGN.CENTER)
    add_text(slide, "run_004 主结果", emu(7.71), emu(3.58), emu(2.65), emu(0.58), size=10.5, color=TEXT,
             align=PP_ALIGN.CENTER)
    add_text(slide, "Selenium 路径：首页 -> 商品详情 -> 加购 -> 下单成功", emu(1.75), emu(4.82), emu(8.3), emu(0.26),
             size=11, color=PURPLE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "run_004 JMeter：sample_count=48076 | avg_latency=93.13ms | qps=53.4178 | error_rate=0.0%",
             emu(1.75), emu(5.30), emu(8.3), emu(0.32), size=10.5, color=TEXT, align=PP_ALIGN.CENTER)
    add_footer_page_no(slide, 6)


def fill_matrix(slide):
    clear_text(slide)
    add_title(slide, "四轮实验矩阵", "前三轮验证链路，run_004 作为正式展示 case")
    cols = [
        ("run_001", "paymentservice\npod-kill\n8s / 4s"),
        ("run_002", "checkoutservice\npod-kill\n9s / 3s"),
        ("run_003", "recommendationservice\nnetwork-delay\n240s / 120s"),
        ("run_004", "recommendationservice\nnetwork-delay\n900s / 300s"),
    ]
    xs = [1.38, 3.92, 6.42, 8.94]
    for x, (title, body) in zip(xs, cols):
        add_text(slide, title, emu(x), emu(1.64), emu(1.55), emu(0.34), size=12.5, color=WHITE,
                 bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, body, emu(x - 0.08), emu(3.88), emu(1.8), emu(0.82), size=10.2, color=TEXT,
                 align=PP_ALIGN.CENTER, line_spacing=1.08)
    add_text(slide, "run_004 具备最长总窗口、最长故障窗口和最完整的故障前/中/后趋势，因此作为答辩主案例。",
             emu(1.45), emu(5.44), emu(9.0), emu(0.38), size=12, color=PURPLE, bold=True,
             align=PP_ALIGN.CENTER)
    add_footer_page_no(slide, 7)


def fill_run004(slide):
    clear_text(slide)
    add_title(slide, "主结果：run_004 正式版 case", "recommendationservice network-delay 的单 case 结果")
    add_rect(slide, emu(0.65), emu(1.10), emu(11.05), emu(4.82), fill=WHITE)
    add_picture_contain(slide, IMG_RUN004, emu(0.78), emu(1.18), emu(7.15), emu(4.18))
    add_rect(slide, emu(8.10), emu(1.52), emu(3.25), emu(2.70), fill=PURPLE)
    add_text(slide, "关键指标", emu(8.45), emu(1.70), emu(2.45), emu(0.32), size=18, color=WHITE,
             bold=True, align=PP_ALIGN.CENTER)
    add_bullets(slide, [
        "总窗口：900 秒",
        "故障窗口：300 秒",
        "故障桶：11 个",
        "宽表规模：32 行",
        "KPIRoot hit@10：1.0",
    ], emu(8.28), emu(2.18), emu(2.75), emu(1.75), size=12, color=WHITE)
    add_text(slide, "展示价值：覆盖故障前、故障中与恢复后，是当前最稳定的正式主案例。",
             emu(1.55), emu(5.48), emu(8.95), emu(0.38), size=12, color=TEXT, align=PP_ALIGN.CENTER)
    add_footer_page_no(slide, 8)


def fill_algorithm_steps(slide):
    clear_text(slide)
    add_title(slide, "算法复现流程", "从宽表输入到异常检测和根因定位输出")
    steps = [
        ("步骤 01", "宽表构建", "Prometheus KPI、JMeter 指标和故障标签按 30s 粒度对齐。"),
        ("步骤 02", "CATCH 输入", "将宽表转为多变量时间序列，输出异常检测分数。"),
        ("步骤 03", "KPIRoot 输入", "生成候选服务 KPI 与告警 KPI，构建根因定位 case。"),
        ("步骤 04", "结果汇总", "生成单 case 图和跨 run 总图，用于报告与答辩。"),
    ]
    xs = [0.92, 3.95, 6.88, 9.78]
    for x, (step, title, body) in zip(xs, steps):
        add_text(slide, step, emu(x + 0.55), emu(2.34), emu(1.25), emu(0.28), size=10.5, color=WHITE,
                 bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, title, emu(x), emu(3.70), emu(2.2), emu(0.28), size=12, color=PURPLE,
                 bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, body, emu(x - 0.02), emu(4.56), emu(2.25), emu(0.76), size=9.2, color=TEXT,
                 align=PP_ALIGN.CENTER, line_spacing=1.08)
    add_footer_page_no(slide, 9)


def fill_cross_result(slide):
    clear_text(slide)
    add_title(slide, "跨 4 个 run 的统一结果", "不同故障类型和窗口长度下的算法表现对比")
    add_picture_contain(slide, IMG_SUMMARY, emu(0.82), emu(1.35), emu(5.1), emu(3.75))
    notes = [
        ("JMeter 稳定性", "四个 run 的错误率均为 0.0%，说明业务路径在实验窗口内稳定。"),
        ("CATCH 观察", "pod-kill 更容易被检测；network-delay 属于更难的性能退化场景。"),
        ("KPIRoot 命中", "四个 run 中真实故障服务均进入候选集合，hit@10 = 1.0。"),
        ("展示选择", "run_004 时间窗更完整，最适合作为正式展示主案例。"),
    ]
    tops = [1.55, 2.78, 4.02, 5.25]
    for top, (title, body) in zip(tops, notes):
        add_text(slide, title, emu(7.1), emu(top), emu(1.55), emu(0.22), size=11.5, color=PURPLE, bold=True)
        add_text(slide, body, emu(8.55), emu(top), emu(2.85), emu(0.48), size=9.4, color=TEXT, line_spacing=1.05)
    add_footer_page_no(slide, 10)


def fill_engineering(slide):
    clear_text(slide)
    add_title(slide, "工程问题与解决办法", "展示环境恢复、数据处理和稳定性保障工作")
    cards = [
        (2.65, 2.20, "环境部署", "Minikube 镜像慢\n预拉取 kicbase 并使用 Docker driver"),
        (8.18, 2.20, "访问入口", "Windows 外部 IP 不稳定\n统一使用 kubectl port-forward"),
        (2.65, 4.42, "数据处理", "ISO8601 时间格式混合\n脚本兼容解析并对齐标签"),
        (8.18, 4.42, "故障清理", "NetworkChaos finalizer 残留\n清空 finalizers 后删除资源"),
    ]
    for x, y, title, body in cards:
        add_text(slide, title, emu(x), emu(y), emu(1.95), emu(0.25), size=11.2, color=WHITE,
                 bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, body, emu(x - 0.08), emu(y + 0.46), emu(2.35), emu(0.58), size=9.6,
                 color=TEXT, align=PP_ALIGN.CENTER, line_spacing=1.08)
    add_text(slide, "工程闭环", emu(5.55), emu(3.45), emu(1.32), emu(0.48), size=14, color=WHITE,
             bold=True, align=PP_ALIGN.CENTER, line_spacing=1.0)
    add_footer_page_no(slide, 11)


def fill_summary(slide):
    clear_text(slide)
    add_title(slide, "总结与后续工作", "当前已经完成可展示、可解释的微服务 AIOps 实验闭环")
    add_text(slide, "已完成", emu(1.58), emu(2.38), emu(1.3), emu(0.56), size=16, color=WHITE,
             bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "部署、业务验证、监控采集、故障注入、CATCH 异常检测与 KPIRoot 根因定位均已打通。\n核心证据包括 run_004 主图、跨 4 个 run 总图和 Selenium 下单成功证据。",
             emu(4.08), emu(2.38), emu(6.35), emu(0.86), size=12, color=WHITE, line_spacing=1.12)
    add_text(slide, "后续工作", emu(9.82), emu(4.42), emu(1.3), emu(0.56), size=15, color=WHITE,
             bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "继续补充 paymentservice / checkoutservice 的 CPU stress 对照实验，增强故障类型多样性与结论说服力。",
             emu(1.55), emu(4.56), emu(6.65), emu(0.56), size=12, color=TEXT, line_spacing=1.12)
    add_footer_page_no(slide, 12)


def fill_thanks(slide):
    clear_text(slide)
    add_text(slide, "感谢您的观看！", emu(2.35), emu(2.92), emu(8.8), emu(0.82),
             size=36, color=WHITE, bold=True, align=PP_ALIGN.CENTER, font=FONT_BOLD)
    add_text(slide, "基于 Online Boutique 的异常检测与根因定位闭环复现",
             emu(2.65), emu(4.15), emu(8.2), emu(0.32), size=14, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, "汇报人：王熙康    学号：2312124",
             emu(3.3), emu(4.78), emu(6.7), emu(0.28), size=11.5, color=WHITE, align=PP_ALIGN.CENTER)


def build_deck():
    prs = Presentation(str(TEMPLATE))
    keep_only_selected_slides(prs, SELECTED_TEMPLATE_SLIDES)

    fillers = [
        fill_cover,
        fill_toc,
        fill_section,
        fill_system_intro,
        fill_closed_loop,
        fill_test_evidence,
        fill_matrix,
        fill_run004,
        fill_algorithm_steps,
        fill_cross_result,
        fill_engineering,
        fill_summary,
        fill_thanks,
    ]
    for slide, filler in zip(prs.slides, fillers):
        filler(slide)

    prs.save(str(OUTPUT))
    print(OUTPUT)


if __name__ == "__main__":
    build_deck()